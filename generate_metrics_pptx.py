"""
Generate PowerPoint slides explaining RL Training Metrics
Output: rl_metrics_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLORS = {
    'bg':     RGBColor(0x0a, 0x0e, 0x14),
    'panel':  RGBColor(0x18, 0x18, 0x1b),
    'panel2': RGBColor(0x27, 0x27, 0x2a),
    'text':   RGBColor(0xe4, 0xe4, 0xe7),
    'muted':  RGBColor(0x71, 0x71, 0x7a),
    'green':  RGBColor(0x34, 0xd3, 0x99),
    'red':    RGBColor(0xef, 0x44, 0x44),
    'yellow': RGBColor(0xfb, 0xbf, 0x24),
    'blue':   RGBColor(0x60, 0xa5, 0xfa),
    'purple': RGBColor(0xa7, 0x8b, 0xfa),
    'orange': RGBColor(0xfb, 0x92, 0x3c),
    'pink':   RGBColor(0xf4, 0x72, 0xb6),
    'white':  RGBColor(0xff, 0xff, 0xff),
}


def set_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or COLORS['bg']


def add_text(slide, left, top, width, height, text,
              size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
              font_name=None):
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return tx


def add_card(slide, left, top, width, height, accent_color=None,
              fill_color=None):
    """Add a styled card (rounded rectangle)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or COLORS['panel']
    if accent_color:
        shape.line.color.rgb = accent_color
        shape.line.width = Pt(2.5)
    else:
        shape.line.color.rgb = COLORS['muted']
        shape.line.width = Pt(0.5)
    return shape


def add_metric_card(slide, left, top, width, height, title, value,
                     status_text="", color=COLORS['blue']):
    """Card showing a metric with status"""
    card = add_card(slide, left, top, width, height, accent_color=color)
    tf = card.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(13)
    r1.font.color.rgb = COLORS['muted']

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.color.rgb = color
    r2.font.name = 'Consolas'

    if status_text:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = status_text
        r3.font.size = Pt(11)
        r3.font.color.rgb = COLORS['text']


# ============================================================
# Build Presentation
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 1, 1.8, 11.3, 1.5, "📊 RL Training Metrics",
         size=58, bold=True, color=COLORS['blue'], align=PP_ALIGN.CENTER)
add_text(s, 1, 3.2, 11.3, 0.7,
         "อ่านค่าระหว่างเทรน — รู้ทันปัญหา ก่อนเสียเวลา 30 นาที",
         size=22, color=COLORS['text'], align=PP_ALIGN.CENTER)
add_text(s, 1, 4.0, 11.3, 0.5,
         "Healthy Metrics สำหรับ PPO + Stable-Baselines3",
         size=16, color=COLORS['muted'], align=PP_ALIGN.CENTER)

add_text(s, 1, 5.2, 11.3, 0.8, "🩺 📈 🎯",
         size=58, align=PP_ALIGN.CENTER)

add_text(s, 1, 6.6, 11.3, 0.4,
         "Course Slides · 14 slides",
         size=12, color=COLORS['muted'], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: WHY METRICS MATTER
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🤔 ทำไมต้องดู Metrics?",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.4, 12, 0.5,
         "เทรน RL ใช้เวลา 30 นาที - 3 ชั่วโมง — ถ้าเข้าใจ metrics จะรู้ว่า:",
         size=16, color=COLORS['muted'])

scenarios = [
    ("✅", "Training กำลังไปได้ดี",
     "→ ปล่อยให้รันต่อไปจนจบ", COLORS['green']),
    ("⚠️", "มีอาการแปลกๆ",
     "→ หยุดทันที อย่าเสียเวลา ปรับแล้วเริ่มใหม่", COLORS['yellow']),
    ("❌", "Training พังแล้ว (collapse)",
     "→ Stop + Debug + Restart", COLORS['red']),
]

for i, (icon, title, action, color) in enumerate(scenarios):
    y = 2.4 + i * 1.4
    add_card(s, 0.5, y, 12.3, 1.2, accent_color=color)
    add_text(s, 0.7, y + 0.15, 0.8, 0.8, icon,
             size=36, align=PP_ALIGN.CENTER)
    add_text(s, 1.6, y + 0.15, 11, 0.4, title,
             size=18, bold=True, color=color)
    add_text(s, 1.6, y + 0.65, 11, 0.4, action,
             size=14, color=COLORS['text'])

add_text(s, 0.5, 6.8, 12, 0.4,
         '"Stop early, save time" — กฎข้อแรกของ ML Engineer',
         size=14, color=COLORS['yellow'], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: 6 KEY METRICS OVERVIEW
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📋 6 Metrics สำคัญ",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ดูตัวเหล่านี้ในระหว่างเทรน — แต่ละตัวบอกคนละเรื่อง",
         size=15, color=COLORS['muted'])

metrics_overview = [
    ("ep_rew_mean",    "💰", "Reward เฉลี่ย", "บอก: model เรียนได้ผลไหม", COLORS['green']),
    ("approx_kl",      "📐", "Policy stability", "บอก: เปลี่ยนเร็วไปไหม", COLORS['blue']),
    ("clip_fraction",  "✂️", "How much clipping", "บอก: stable vs aggressive", COLORS['purple']),
    ("explained_var",  "🎯", "Critic quality", "บอก: ทาย value ได้แม่นไหม", COLORS['pink']),
    ("entropy_loss",   "🔍", "Exploration level", "บอก: explore vs exploit", COLORS['orange']),
    ("value_loss",     "📉", "Critic learning", "บอก: critic ฝึกได้ดีไหม", COLORS['yellow']),
]

for i, (name, icon, title, desc, color) in enumerate(metrics_overview):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.2
    y = 2.0 + row * 2.4

    card = add_card(s, x, y, 4, 2.2, accent_color=color)
    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = icon
    r1.font.size = Pt(36)

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = name
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = color
    r2.font.name = 'Consolas'

    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = title
    r3.font.size = Pt(13)
    r3.font.color.rgb = COLORS['text']

    p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run(); r4.text = desc
    r4.font.size = Pt(11)
    r4.font.color.rgb = COLORS['muted']


# ============================================================
# SLIDE 4: ep_rew_mean (most important!)
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "💰 ep_rew_mean",
         size=32, bold=True, color=COLORS['green'])
add_text(s, 0.5, 1.2, 12, 0.5, "Reward เฉลี่ยต่อ episode — ตัวสำคัญที่สุด",
         size=16, color=COLORS['muted'])

add_text(s, 0.5, 2.0, 12, 0.5,
         "ทำหน้าที่: บอกว่า agent เรียนรู้ได้ผลไหม",
         size=14, color=COLORS['text'], bold=True)

# Example trends
add_text(s, 0.5, 2.7, 12, 0.4,
         "Pattern ที่ควรเห็น:",
         size=14, color=COLORS['yellow'], bold=True)

trends = [
    ("📈 ขึ้นเรื่อยๆ",
     "Step 0:  -0.5\nStep 50k: 0.5\nStep 100k: 1.5\nStep 200k: 2.5",
     "✅ กำลังเรียนรู้", COLORS['green']),
    ("➡️ ทรงตัว",
     "Step 0:  -0.1\nStep 50k: 0.0\nStep 100k: 0.0\nStep 200k: 0.0",
     "🟡 ไม่เรียน → check approx_kl", COLORS['yellow']),
    ("📉 ลดลง / ลบ",
     "Step 0:  -0.5\nStep 50k: -1.0\nStep 100k: -2.0\nStep 200k: -5.0",
     "❌ Collapse → Stop!", COLORS['red']),
]

for i, (title, data, verdict, color) in enumerate(trends):
    x = 0.5 + i * 4.2
    card = add_card(s, x, 3.2, 4, 3.5, accent_color=color)
    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = color

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = data
    r2.font.size = Pt(12); r2.font.name = 'Consolas'
    r2.font.color.rgb = COLORS['text']

    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = verdict
    r3.font.size = Pt(13); r3.font.color.rgb = color


# ============================================================
# SLIDE 5: approx_kl
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📐 approx_kl",
         size=32, bold=True, color=COLORS['blue'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "KL Divergence — Policy ใหม่ต่างจาก Policy เก่าแค่ไหน",
         size=16, color=COLORS['muted'])

# Concept
add_text(s, 0.5, 2.0, 12, 0.5,
         '"Proximal" ใน PPO = ห้าม Policy เปลี่ยนเร็วเกิน → approx_kl คือตัววัด',
         size=14, color=COLORS['text'])

# Healthy ranges
add_text(s, 0.5, 2.8, 12, 0.4, "🩺 Healthy Ranges:",
         size=15, color=COLORS['yellow'], bold=True)

ranges = [
    ("0.005-0.02",  "🟢 ดีมาก", "Smooth, stable learning",         COLORS['green']),
    ("0.02-0.05",   "🟢 ดี",    "Normal",                          COLORS['green']),
    ("0.05-0.10",   "🟡 พอใช้", "Watch carefully",                 COLORS['yellow']),
    ("0.10-0.30",   "🔴 แย่",   "Unstable",                        COLORS['red']),
    ("> 0.30",      "💥 พัง",   "Policy collapse",                 COLORS['red']),
]

for i, (val, status, note, color) in enumerate(ranges):
    y = 3.3 + i * 0.5
    add_card(s, 0.5, y, 12.3, 0.45, accent_color=color)
    add_text(s, 0.7, y + 0.05, 2.5, 0.35, val,
             size=14, bold=True, color=color, font_name='Consolas')
    add_text(s, 3.5, y + 0.05, 2.5, 0.35, status,
             size=14, color=COLORS['text'])
    add_text(s, 6.5, y + 0.05, 6, 0.35, note,
             size=13, color=COLORS['muted'])

add_text(s, 0.5, 6.0, 12, 0.4,
         '⚠️  ของคุณตอนนี้ approx_kl = 0.51 → 10× เกิน healthy',
         size=14, color=COLORS['red'], bold=True)
add_text(s, 0.5, 6.5, 12, 0.4,
         '   วิธีแก้: ลด window_size, ลด learning_rate, ลด clip_range',
         size=13, color=COLORS['muted'])


# ============================================================
# SLIDE 6: clip_fraction
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "✂️ clip_fraction",
         size=32, bold=True, color=COLORS['purple'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "% ของ updates ที่โดน clip (ห้ามเปลี่ยนเกิน 20%)",
         size=16, color=COLORS['muted'])

# Visual analogy
analogy_card = add_card(s, 0.5, 2.0, 12.3, 2.5, accent_color=COLORS['purple'])
tf = analogy_card.text_frame
tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "💡 เปรียบเหมือนการบังคับเลี้ยวรถ:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['purple']

for line in [
    "",
    "  clip_range = 0.2  →  หมุนพวงมาลัยได้สูงสุด 20% ต่อครั้ง",
    "  clip_fraction = 0.05  →  5% ของครั้งที่เลี้ยวเต็ม (= 95% เลี้ยวเบาๆ)  ✅",
    "  clip_fraction = 0.50  →  50% ของครั้งที่เลี้ยวสุด (= ขับห่ามมาก!)   🔴",
]:
    pp = tf.add_paragraph()
    rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(13); rr.font.name = 'Consolas'
    rr.font.color.rgb = COLORS['text']

# Healthy ranges
add_text(s, 0.5, 4.7, 12, 0.4, "🩺 Healthy Ranges:",
         size=15, color=COLORS['yellow'], bold=True)

cf_ranges = [
    ("0.0 - 0.1", "🟢 ดีมาก", "Stable", COLORS['green']),
    ("0.1 - 0.2", "🟢 ดี",    "Normal", COLORS['green']),
    ("0.2 - 0.3", "🟡 กลาง",   "Borderline", COLORS['yellow']),
    ("> 0.3",     "🔴 แย่",    "Policy unstable", COLORS['red']),
]

for i, (val, status, note, color) in enumerate(cf_ranges):
    y = 5.2 + i * 0.45
    add_card(s, 0.5, y, 12.3, 0.4, accent_color=color)
    add_text(s, 0.7, y + 0.03, 2.5, 0.3, val,
             size=14, bold=True, color=color, font_name='Consolas')
    add_text(s, 3.5, y + 0.03, 2.5, 0.3, status,
             size=14, color=COLORS['text'])
    add_text(s, 6.5, y + 0.03, 6, 0.3, note,
             size=13, color=COLORS['muted'])


# ============================================================
# SLIDE 7: explained_variance
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🎯 explained_variance",
         size=32, bold=True, color=COLORS['pink'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "Critic ทาย expected reward ได้แม่นแค่ไหน",
         size=16, color=COLORS['muted'])

add_text(s, 0.5, 2.0, 12, 0.5,
         "Critic = NN ตัวที่ 2 ใน PPO — ทำนาย 'state นี้คาดว่าได้ reward เท่าไหร่'",
         size=14, color=COLORS['text'])

# Healthy ranges
add_text(s, 0.5, 2.8, 12, 0.4, "🩺 Healthy Ranges:",
         size=15, color=COLORS['yellow'], bold=True)

ev_ranges = [
    ("0.8 - 1.0", "🟢 ดีมาก", "Critic ทายแม่น 80%+", COLORS['green']),
    ("0.5 - 0.8", "🟢 ดี",    "Useful predictions", COLORS['green']),
    ("0.2 - 0.5", "🟡 พอใช้", "Marginal", COLORS['yellow']),
    ("0 - 0.2",   "🔴 แย่",   "Critic ใช้ไม่ได้", COLORS['red']),
    ("< 0",       "💥 พัง",   "ทายแย่กว่าค่าเฉลี่ย", COLORS['red']),
]

for i, (val, status, note, color) in enumerate(ev_ranges):
    y = 3.3 + i * 0.5
    add_card(s, 0.5, y, 12.3, 0.45, accent_color=color)
    add_text(s, 0.7, y + 0.05, 2.5, 0.35, val,
             size=14, bold=True, color=color, font_name='Consolas')
    add_text(s, 3.5, y + 0.05, 2.5, 0.35, status,
             size=14, color=COLORS['text'])
    add_text(s, 6.5, y + 0.05, 6, 0.35, note,
             size=13, color=COLORS['muted'])

add_text(s, 0.5, 6.0, 12, 0.4,
         "💡 Pattern: เริ่มต้น 0 → เพิ่มขึ้นเรื่อยๆ → 0.7+ ตอนจบ",
         size=14, color=COLORS['blue'], bold=True)
add_text(s, 0.5, 6.5, 12, 0.4,
         "    ถ้า ev < 0 = critic ทำลาย training → อาจต้อง restart",
         size=13, color=COLORS['muted'])


# ============================================================
# SLIDE 8: entropy_loss
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🔍 entropy_loss",
         size=32, bold=True, color=COLORS['orange'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "ระดับการสำรวจ (exploration) ของ agent",
         size=16, color=COLORS['muted'])

# Concept
add_text(s, 0.5, 2.0, 12, 0.5,
         "entropy สูง = สำรวจหลาย actions  ·  entropy ต่ำ = มั่นใจ ใช้ action เดิม",
         size=14, color=COLORS['text'])

add_text(s, 0.5, 2.6, 12, 0.4,
         "หมายเหตุ: รายงานเป็น negative loss → ค่าลบมากกว่า = exploration มากกว่า",
         size=13, color=COLORS['muted'])

# Phases
add_text(s, 0.5, 3.3, 12, 0.4, "📅 Healthy Pattern ตลอดการเทรน:",
         size=15, color=COLORS['yellow'], bold=True)

phases = [
    ("Early\n0-30%",   "-1.2 to -1.0",  "🌱 สำรวจมาก", COLORS['green']),
    ("Mid\n30-70%",    "-0.8 to -0.5",  "⚖️ Balance",   COLORS['blue']),
    ("Late\n70-100%",  "-0.5 to -0.2",  "🎯 Exploit", COLORS['purple']),
]

for i, (phase, val, note, color) in enumerate(phases):
    x = 0.5 + i * 4.2
    card = add_card(s, x, 4.0, 4, 2.5, accent_color=color)
    tf = card.text_frame
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.3)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = phase
    r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = color

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = val
    r2.font.size = Pt(18); r2.font.name = 'Consolas'
    r2.font.color.rgb = COLORS['text']

    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = note
    r3.font.size = Pt(14); r3.font.color.rgb = color

add_text(s, 0.5, 6.8, 12, 0.4,
         "🎛️ ปรับด้วย ent_coef (default 0.01) — ถ้าน้อยเกิน policy collapse เร็ว",
         size=12, color=COLORS['muted'])


# ============================================================
# SLIDE 9: value_loss
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📉 value_loss",
         size=32, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "Critic Training Loss (Mean Squared Error)",
         size=16, color=COLORS['muted'])

add_text(s, 0.5, 2.0, 12, 0.5,
         "ความผิดพลาดของ Critic เมื่อทาย expected reward",
         size=14, color=COLORS['text'])

add_text(s, 0.5, 2.7, 12, 0.4,
         "value_loss = (V_predicted − V_actual)²  → ยิ่งต่ำยิ่งดี",
         size=13, color=COLORS['blue'], font_name='Consolas')

# Ranges
add_text(s, 0.5, 3.5, 12, 0.4, "🩺 Healthy Ranges:",
         size=15, color=COLORS['yellow'], bold=True)

vl_ranges = [
    ("< 0.01",     "🟢 ดี",     "Critic เรียนได้ดี", COLORS['green']),
    ("0.01 - 0.1", "🟡 กลาง",  "OK", COLORS['yellow']),
    ("> 0.1",      "🔴 แย่",    "Critic struggling", COLORS['red']),
]

for i, (val, status, note, color) in enumerate(vl_ranges):
    y = 4.0 + i * 0.55
    add_card(s, 0.5, y, 12.3, 0.5, accent_color=color)
    add_text(s, 0.7, y + 0.06, 2.5, 0.35, val,
             size=14, bold=True, color=color, font_name='Consolas')
    add_text(s, 3.5, y + 0.06, 2.5, 0.35, status,
             size=14, color=COLORS['text'])
    add_text(s, 6.5, y + 0.06, 6, 0.35, note,
             size=13, color=COLORS['muted'])

add_text(s, 0.5, 6.2, 12, 0.4,
         "✅ ของคุณตอนนี้: 0.000393 → ดีเยี่ยม! Critic ทำงานปกติ",
         size=14, color=COLORS['green'], bold=True)


# ============================================================
# SLIDE 10: HEALTHY METRICS TABLE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📋 Healthy Metrics Table",
         size=32, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.5, "ตารางสรุปค่าที่ปลอดภัย",
         size=15, color=COLORS['muted'])

# Header
y0 = 2.0
add_card(s, 0.5, y0, 12.3, 0.5, fill_color=COLORS['panel2'])
headers = [("Metric", 0.5, 2.5), ("🟢 Healthy", 3.0, 2.5),
            ("🟡 Warning", 5.5, 2.5), ("🔴 Bad", 8.0, 2.5),
            ("ของคุณ (ปัจจุบัน)", 10.5, 2.5)]
for h, x, w in headers:
    add_text(s, x + 0.1, y0 + 0.1, w, 0.3, h,
             size=13, bold=True, color=COLORS['yellow'])

# Rows
metrics_data = [
    ("approx_kl",      "0.005-0.05",  "0.05-0.15",  "> 0.15",   "0.51 ❌", COLORS['red']),
    ("clip_fraction",  "< 0.2",       "0.2-0.3",    "> 0.3",    "0.52 ❌", COLORS['red']),
    ("explained_var",  "> 0.5",       "0.2-0.5",    "< 0.2",    "0.50 🟡", COLORS['yellow']),
    ("entropy_loss",   "-1 to -0.3",  "-0.3 to -0.1", "> -0.1", "-0.55 🟡", COLORS['yellow']),
    ("value_loss",     "< 0.01",      "0.01-0.1",   "> 0.1",    "0.0004 🟢", COLORS['green']),
]

for i, (m, h, w, b, your, your_color) in enumerate(metrics_data):
    y = y0 + 0.55 + i * 0.5
    bg_color = COLORS['panel'] if i % 2 == 0 else COLORS['panel2']
    add_card(s, 0.5, y, 12.3, 0.45, fill_color=bg_color)

    add_text(s, 0.7, y + 0.05, 2.4, 0.35, m,
             size=12, bold=True, color=COLORS['blue'], font_name='Consolas')
    add_text(s, 3.1, y + 0.05, 2.4, 0.35, h,
             size=12, color=COLORS['green'], font_name='Consolas')
    add_text(s, 5.6, y + 0.05, 2.4, 0.35, w,
             size=12, color=COLORS['yellow'], font_name='Consolas')
    add_text(s, 8.1, y + 0.05, 2.4, 0.35, b,
             size=12, color=COLORS['red'], font_name='Consolas')
    add_text(s, 10.6, y + 0.05, 2.4, 0.35, your,
             size=12, bold=True, color=your_color, font_name='Consolas')

add_text(s, 0.5, 6.0, 12, 0.4,
         "🩺 Diagnosis: Window=50 ใหญ่เกิน → policy unstable",
         size=15, color=COLORS['red'], bold=True)
add_text(s, 0.5, 6.5, 12, 0.4,
         "💊 แก้: ลด window=10 (default) → metrics ทุกตัวกลับมา healthy",
         size=14, color=COLORS['green'])


# ============================================================
# SLIDE 11: HEALTHY EXAMPLE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "✅ Healthy Training Example",
         size=32, bold=True, color=COLORS['green'])
add_text(s, 0.5, 1.2, 12, 0.5, "rl_v3 หลังเทรน 200k steps — ค่าทุกตัวอยู่ใน healthy range",
         size=15, color=COLORS['muted'])

# Stats grid
healthy_stats = [
    ("ep_rew_mean",        "+5.24",    "✅ positive", COLORS['green']),
    ("approx_kl",          "0.012",    "✅ smooth",   COLORS['green']),
    ("clip_fraction",      "0.08",     "✅ low",      COLORS['green']),
    ("explained_variance", "0.74",     "✅ good",     COLORS['green']),
    ("entropy_loss",       "-1.35",    "✅ exploring", COLORS['green']),
    ("value_loss",         "0.0008",   "✅ excellent", COLORS['green']),
]

for i, (label, value, status, color) in enumerate(healthy_stats):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.2
    y = 2.2 + row * 1.7
    add_metric_card(s, x, y, 4, 1.5, label, value, status, color)

# Result
add_card(s, 0.5, 5.8, 12.3, 1.3, accent_color=COLORS['green'])
tf = prs.slides[-1].shapes[-1].text_frame
tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "🎯 Backtest Result:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['green']

p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "  PF = 1.02 · WR = 52.4% · Return = +1.92% · Max DD = -8.6%"
r2.font.size = Pt(14); r2.font.color.rgb = COLORS['text']
r2.font.name = 'Consolas'


# ============================================================
# SLIDE 12: UNHEALTHY EXAMPLE (your current)
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "⚠️  Unhealthy Training (ของคุณตอนนี้)",
         size=30, bold=True, color=COLORS['red'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "rl_v3_new กับ window=50 — มีอาการแปลกๆ",
         size=15, color=COLORS['muted'])

unhealthy_stats = [
    ("approx_kl",          "0.51",    "🔴 10× เกิน",     COLORS['red']),
    ("clip_fraction",      "0.52",    "🔴 unstable",   COLORS['red']),
    ("explained_variance", "0.50",    "🟡 borderline", COLORS['yellow']),
    ("entropy_loss",       "-0.55",   "🟡 ลดเร็วไป",     COLORS['yellow']),
    ("value_loss",         "0.0004",  "🟢 OK",         COLORS['green']),
    ("n_updates",          "170",     "📊 18% done",   COLORS['blue']),
]

for i, (label, value, status, color) in enumerate(unhealthy_stats):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.2
    y = 2.2 + row * 1.7
    add_metric_card(s, x, y, 4, 1.5, label, value, status, color)

# Diagnosis card
add_card(s, 0.5, 5.8, 12.3, 1.3, accent_color=COLORS['red'])
diag_card = prs.slides[-1].shapes[-1]
tf = diag_card.text_frame
tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "🩺 Root Cause:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['red']

p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "  Window=50 ใหญ่เกิน → state ใหญ่ → NN ปรับตัวยาก → policy กระโดด"
r2.font.size = Pt(14); r2.font.color.rgb = COLORS['text']

p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = "💊 Fix: เปลี่ยน window=10 (default) → restart training"
r3.font.size = Pt(14); r3.font.bold = True; r3.font.color.rgb = COLORS['green']


# ============================================================
# SLIDE: HOW TO FIX EACH METRIC
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🔧 How to Fix Each Metric",
         size=32, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "ตัวที่แย่แก้ได้เสมอ — ลองตามลำดับด้านล่าง",
         size=15, color=COLORS['muted'])

# Header
y0 = 2.0
add_card(s, 0.5, y0, 12.3, 0.5, fill_color=COLORS['panel2'])
add_text(s, 0.7, y0 + 0.1, 3, 0.3, "Problem",
         size=13, bold=True, color=COLORS['yellow'])
add_text(s, 4.0, y0 + 0.1, 4.5, 0.3, "Primary Fix ⭐",
         size=13, bold=True, color=COLORS['yellow'])
add_text(s, 9.0, y0 + 0.1, 4, 0.3, "Backup Fix",
         size=13, bold=True, color=COLORS['yellow'])

fixes = [
    ("approx_kl > 0.05",      "↓ learning_rate (3e-4 → 1e-4)",  "↓ clip_range, ↓ n_epochs", COLORS['blue']),
    ("clip_fraction > 0.3",   "↓ learning_rate (3e-4 → 1e-4)",  "↑ n_steps, ↑ clip_range",  COLORS['purple']),
    ("explained_var < 0.5",   "↑ Critic NN size",                "↑ gae_lambda, ↑ vf_coef",  COLORS['pink']),
    ("entropy → 0 เร็วไป",     "↑ ent_coef (0.01 → 0.05)",       "↓ learning_rate",          COLORS['orange']),
    ("value_loss > 0.1",      "Reward normalization",            "↑ Critic NN size",         COLORS['yellow']),
    ("ep_rew_mean stuck",     "Reward shaping",                   "↑ features, ↑ ent_coef",   COLORS['green']),
]

for i, (problem, primary, backup, color) in enumerate(fixes):
    y = y0 + 0.55 + i * 0.7
    bg = COLORS['panel'] if i % 2 == 0 else COLORS['panel2']
    add_card(s, 0.5, y, 12.3, 0.65, fill_color=bg)
    add_text(s, 0.7, y + 0.15, 3.2, 0.4, problem,
             size=12, color=color, bold=True, font_name='Consolas')
    add_text(s, 4.0, y + 0.15, 4.8, 0.4, primary,
             size=12, color=COLORS['green'])
    add_text(s, 9.0, y + 0.15, 4, 0.4, backup,
             size=11, color=COLORS['muted'])

add_text(s, 0.5, 6.8, 12, 0.4,
         "💡 Tip: ลอง Primary fix ก่อน · Train 50k steps ดูผล · ถ้ายังแย่ลอง Backup",
         size=13, color=COLORS['blue'], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE: clip_fraction Deep Dive (counter-intuitive)
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "✂️ clip_fraction — Counter-intuitive Fix",
         size=28, bold=True, color=COLORS['purple'])
add_text(s, 0.5, 1.2, 12, 0.5,
         "หลายคนคิดว่า 'ลด clip_range' จะแก้ clip_fraction → จริงๆ ตรงข้าม!",
         size=15, color=COLORS['muted'])

# Wrong intuition
wrong = add_card(s, 0.5, 2.0, 6, 4.5, accent_color=COLORS['red'])
tf = wrong.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "❌ คิดผิด"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = COLORS['red']

for line in [
    "",
    "'clip_fraction สูง = clip บ่อย",
    "→ ลด clip_range เพื่อ clip น้อยลง'",
    "",
    "ผล:",
    "  clip_range = 0.2 → cf = 0.50",
    "  clip_range = 0.1 → cf = 0.70 ❌",
    "  clip_range = 0.05 → cf = 0.85 ❌",
    "",
    "เพราะ window แคบลง → policy",
    "ติด clip บ่อยขึ้น (ตรงข้ามกับที่คิด)",
]:
    pp = tf.add_paragraph()
    rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(13); rr.font.color.rgb = COLORS['text']
    if line.endswith("❌"):
        rr.font.color.rgb = COLORS['red']

# Right approach
right = add_card(s, 6.8, 2.0, 6, 4.5, accent_color=COLORS['green'])
tf = right.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "✅ คิดถูก"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = COLORS['green']

for line in [
    "",
    "'cf สูง = policy พยายามเปลี่ยน",
    " เร็วเกิน → ลด learning_rate'",
    "",
    "ผล (ลด lr 3e-4 → 1e-4):",
    "  cf: 0.50 → 0.18 ✅",
    "  approx_kl: 0.51 → 0.04 ✅",
    "  → policy stable + เรียนรู้ได้",
    "",
    "หลัก: แก้ที่ต้นเหตุ (lr)",
    "  ไม่ใช่ปลายเหตุ (clip_range)",
]:
    pp = tf.add_paragraph()
    rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(13); rr.font.color.rgb = COLORS['text']
    if line.endswith("✅"):
        rr.font.color.rgb = COLORS['green']

add_text(s, 0.5, 6.8, 12, 0.4,
         "🎯 Rule: Symptoms (clip_fraction) ≠ Cause · ดูที่ root cause (learning_rate)",
         size=14, color=COLORS['yellow'], align=PP_ALIGN.CENTER, bold=True)


# ============================================================
# SLIDE 13: DIAGNOSIS FLOWCHART
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🩺 Diagnosis Flowchart",
         size=32, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.5, "เจออาการแบบไหน → ดูที่ metric ไหน → แก้ยังไง",
         size=15, color=COLORS['muted'])

flow_data = [
    ("Reward ไม่ขึ้น",       "approx_kl",  "ปรับ learning_rate (3e-4 → 1e-4)"),
    ("Training crashed",     "approx_kl",  "ลด lr / window / batch_size"),
    ("Trade ทุก bar / ไม่ปิด", "entropy",    "Reward shaping (realized + max_hold)"),
    ("Policy collapse / เลือก action เดียว", "entropy",  "เพิ่ม ent_coef (0.01 → 0.05)"),
    ("Backtest แย่กว่า train", "explained_var", "Overfit → reduce features / window"),
    ("ค่าทุกตัวดูดี แต่ไม่ converge", "ep_rew_mean", "เปลี่ยน reward function"),
]

# Header
y0 = 2.0
add_card(s, 0.5, y0, 12.3, 0.5, fill_color=COLORS['panel2'])
add_text(s, 0.7, y0 + 0.1, 4, 0.3, "อาการ",
         size=13, bold=True, color=COLORS['yellow'])
add_text(s, 5.0, y0 + 0.1, 3, 0.3, "ดูที่ Metric",
         size=13, bold=True, color=COLORS['yellow'])
add_text(s, 8.0, y0 + 0.1, 5, 0.3, "วิธีแก้",
         size=13, bold=True, color=COLORS['yellow'])

for i, (sym, metric, fix) in enumerate(flow_data):
    y = y0 + 0.55 + i * 0.7
    bg = COLORS['panel'] if i % 2 == 0 else COLORS['panel2']
    add_card(s, 0.5, y, 12.3, 0.65, fill_color=bg)

    add_text(s, 0.7, y + 0.15, 4.2, 0.4, sym,
             size=13, color=COLORS['text'])
    add_text(s, 5.0, y + 0.15, 3, 0.4, metric,
             size=13, color=COLORS['blue'], bold=True, font_name='Consolas')
    add_text(s, 8.0, y + 0.15, 5, 0.4, fix,
             size=12, color=COLORS['green'])


# ============================================================
# SLIDE 14: QUICK REFERENCE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "⚡ Quick Reference Cheat Sheet",
         size=32, bold=True, color=COLORS['yellow'])

# Big visual
ref_data = [
    ("ep_rew_mean",       "🔼 ขึ้นเรื่อยๆ = ดี",                         COLORS['green']),
    ("approx_kl",         "0.005-0.05 = healthy   ·   > 0.15 = แย่",     COLORS['blue']),
    ("clip_fraction",     "< 0.2 = healthy   ·   > 0.3 = แย่",            COLORS['purple']),
    ("explained_variance","> 0.5 = healthy   ·   < 0.2 = แย่",            COLORS['pink']),
    ("entropy_loss",      "เริ่ม -1.0 → จบ -0.3   (ลดลงเรื่อยๆ)",         COLORS['orange']),
    ("value_loss",        "< 0.01 = healthy   ·   ลดลงเรื่อยๆ = ดี",       COLORS['yellow']),
]

for i, (metric, desc, color) in enumerate(ref_data):
    y = 1.6 + i * 0.85
    card = add_card(s, 0.5, y, 12.3, 0.75, accent_color=color)
    tf = card.text_frame
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = metric + "    "
    r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = color; r1.font.name = 'Consolas'

    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(14); r2.font.color.rgb = COLORS['text']

# Bottom tip
add_text(s, 0.5, 6.9, 12, 0.4,
         "💡 ดู ep_rew_mean ก่อน — ถ้าขึ้น = OK · ถ้าไม่ → ค่อยดู metrics ตัวอื่น",
         size=14, color=COLORS['yellow'], align=PP_ALIGN.CENTER, bold=True)


# ============================================================
# SLIDE 15: CLOSING
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0, 1.5, 13.3, 1.5, "🎓 รู้ Metrics = ประหยัดเวลา",
         size=48, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

add_text(s, 0, 3.3, 13.3, 0.6,
         "หยุดทันที ก่อนเสีย 30 นาที — Train ใหม่ดีกว่า ทน",
         size=22, color=COLORS['text'], align=PP_ALIGN.CENTER)

add_text(s, 0, 4.2, 13.3, 0.5,
         "Healthy metrics → Healthy model → Profitable trading",
         size=16, color=COLORS['muted'], align=PP_ALIGN.CENTER)

add_text(s, 0, 5.0, 13.3, 1, "🩺 📊 ✅",
         size=58, align=PP_ALIGN.CENTER)

add_text(s, 0, 6.5, 13.3, 0.4,
         "Now: Stop your unhealthy run → Restart with window=10",
         size=12, color=COLORS['green'], align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output = "rl_metrics_slides_v2.pptx"
prs.save(output)
print(f"[OK] Saved: {output}")
print(f"     Slides: {len(prs.slides)}")
