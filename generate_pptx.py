"""
Generate PowerPoint slides automatically
========================================
สร้าง .pptx จาก content ของเรา ใช้ python-pptx

ติดตั้ง:
    pip install python-pptx

Usage:
    python generate_pptx.py
    -> สร้าง rl_course_slides.pptx
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color palette (RGB)
COLORS = {
    'bg':       RGBColor(0x0a, 0x0e, 0x14),
    'panel':    RGBColor(0x18, 0x18, 0x1b),
    'text':     RGBColor(0xe4, 0xe4, 0xe7),
    'muted':    RGBColor(0x71, 0x71, 0x7a),
    'p1':       RGBColor(0x60, 0xa5, 0xfa),  # blue
    'p2':       RGBColor(0x34, 0xd3, 0x99),  # green
    'p3':       RGBColor(0xa7, 0x8b, 0xfa),  # purple
    'p4':       RGBColor(0xfb, 0x92, 0x3c),  # orange
    'yellow':   RGBColor(0xfb, 0xbf, 0x24),
    'red':      RGBColor(0xef, 0x44, 0x44),
    'white':    RGBColor(0xff, 0xff, 0xff),
}


def set_bg(slide, color):
    """Set slide background color"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text,
              size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add a text box"""
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tx


def add_card(slide, left, top, width, height,
              title, content, accent_color=COLORS['yellow']):
    """Add a colored card with title + content"""
    # background panel
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['panel']
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(2)

    # title
    title_tx = shape.text_frame
    title_tx.margin_left = Inches(0.2)
    title_tx.margin_top = Inches(0.15)
    title_tx.word_wrap = True

    p = title_tx.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = accent_color

    # content
    p2 = title_tx.add_paragraph()
    run2 = p2.add_run()
    run2.text = content
    run2.font.size = Pt(11)
    run2.font.color.rgb = COLORS['text']


def add_code(slide, left, top, width, height, code):
    """Add a code block"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x09, 0x09, 0x0b)
    shape.line.color.rgb = COLORS['muted']

    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.name = 'Consolas'
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS['p2']


# ============================================================
# Build slides
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# === SLIDE 1: Title ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 1, 2, 11.3, 1.5, "🤖 RL Trading EA",
         size=60, bold=True, color=COLORS['p1'], align=PP_ALIGN.CENTER)
add_text(s, 1, 3.7, 11.3, 0.6,
         "การพัฒนา Trading System ด้วย Reinforcement Learning",
         size=22, color=COLORS['text'], align=PP_ALIGN.CENTER)
add_text(s, 1, 4.4, 11.3, 0.6,
         "From Data → Train → Deploy → Maintain",
         size=18, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 1, 5.5, 11.3, 1, "🎯 📊 🚀",
         size=60, align=PP_ALIGN.CENTER)
add_text(s, 1, 6.7, 11.3, 0.4,
         "Course Slides · เริ่มจากศูนย์จนใช้งานจริง",
         size=12, color=COLORS['muted'], align=PP_ALIGN.CENTER)

# === SLIDE 2: Roadmap ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🗺️ Roadmap — 4 Phases",
         size=32, bold=True, color=COLORS['yellow'])

phases = [
    ("Phase 1", "📥", "COLLECT", "เก็บข้อมูลจาก MT4", "1-2 สัปดาห์", COLORS['p1']),
    ("Phase 2", "🧠", "TRAIN/TEST", "เทรน RL + validate", "2-4 สัปดาห์", COLORS['p2']),
    ("Phase 3", "🌐", "DEPLOY", "เชื่อม MT5 + เทรดจริง", "1-3 เดือน", COLORS['p3']),
    ("Phase 4", "🔄", "MAINTAIN", "Fine-tune ทุก 3 เดือน", "ต่อเนื่อง", COLORS['p4']),
]
for i, (num, icon, name, desc, time, color) in enumerate(phases):
    left = 0.5 + i * 3.2
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(2),
                                Inches(2.9), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['panel']
    shape.line.color.rgb = color
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.3)
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.add_run().text = num
    p1.runs[0].font.size = Pt(14)
    p1.runs[0].font.color.rgb = COLORS['muted']

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    p2.add_run().text = icon
    p2.runs[0].font.size = Pt(48)

    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    p3.add_run().text = name
    p3.runs[0].font.size = Pt(20)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = color

    p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
    p4.add_run().text = desc
    p4.runs[0].font.size = Pt(12)
    p4.runs[0].font.color.rgb = COLORS['text']

    p5 = tf.add_paragraph(); p5.alignment = PP_ALIGN.CENTER
    p5.add_run().text = time
    p5.runs[0].font.size = Pt(11)
    p5.runs[0].font.color.rgb = COLORS['muted']

add_text(s, 0.5, 6.6, 12, 0.5,
         "หลัก: ML Trading = Engineering ที่ต้องบำรุงรักษา ไม่ใช่ \"magic ทำครั้งเดียวจบ\"",
         size=14, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE: Key Concepts (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "💡 Concepts ที่ต้องรู้ก่อน",
         size=28, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "คำศัพท์สำคัญที่จะเจอตลอดคอร์ส — รู้ไว้ก่อนช่วยให้เข้าใจง่ายขึ้น",
         size=13, color=COLORS['muted'])

concepts = [
    ("🎯 ML (Machine Learning)",
     "คอมพิวเตอร์เรียนรู้จากข้อมูล — ดูประวัติ → ทาย UP/DOWN"),
    ("🤖 RL (Reinforcement Learning)",
     "ML แบบลองทำ + ได้ reward → เรียนเก่งขึ้น (เหมือนสอนสุนัข)"),
    ("🧠 Agent",
     "\"สมอง\" ที่ตัดสินใจ — ในนี้คือ Neural Network"),
    ("🌍 Environment",
     "\"โลก\" ที่ Agent อยู่ — ในนี้คือตลาดจำลอง"),
    ("🎁 Reward",
     "คะแนน บอกว่าทำดีหรือแย่ — ในนี้คือ P&L ของ trade"),
    ("📐 Features",
     "ตัวแปร input ของ model — RSI, EMA, ATR ฯลฯ (ยิ่งดี → model ฉลาดขึ้น)"),
]
for i, (title, content) in enumerate(concepts):
    row = i // 3
    col = i % 3
    add_card(s, 0.5 + col * 4.2, 1.9 + row * 2.5, 4, 2.3,
             title, content, accent_color=COLORS['yellow'])

# === SLIDE 3: Phase 1 Title ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0, 1.5, 13.3, 0.5, "PHASE 1",
         size=18, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 2.2, 13.3, 1.5, "📥 DATA COLLECTION",
         size=60, bold=True, color=COLORS['p1'], align=PP_ALIGN.CENTER)
add_text(s, 0, 4, 13.3, 1.5, "📊", size=80, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.5, 13.3, 0.6,
         "\"Garbage in, garbage out — ข้อมูลคือทุกอย่าง\"",
         size=20, color=COLORS['text'], align=PP_ALIGN.CENTER)

# === SLIDE 4: Phase 1 Features ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "📥 35 Features ที่เก็บ",
         size=28, bold=True, color=COLORS['p1'])

features = [
    ("📈 Base (5)", "RSI 14, EMA 20/50/200, ATR 14"),
    ("🌐 Multi-TF (5)", "RSI H4/D1, Trend H4/D1, Distance EMA200"),
    ("⚡ Momentum (7)", "MACD ×3, Stochastic K/D, CCI, Williams%R"),
    ("🌪️ Volatility (3)", "BB position, ATR ratio, ADX"),
    ("⏰ Time (5)", "Hour, DOW, 3 Sessions"),
    ("📊 Statistical (10)", "Lagged returns, z-score, Sharpe"),
]
for i, (title, content) in enumerate(features):
    row = i // 3
    col = i % 3
    add_card(s, 0.5 + col * 4.2, 1.5 + row * 1.6, 4, 1.4,
             title, content, accent_color=COLORS['p1'])

add_text(s, 0.5, 5, 12, 0.5,
         "= 35 Features × 100,000+ Bars (XAUUSD H1, 2010-2026)",
         size=18, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE: Why 35 features? (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🤔 ทำไมต้อง 35 features?",
         size=28, bold=True, color=COLORS['p1'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "เปรียบเทียบ — ราคาดิบ vs features ที่ \"แปลความหมาย\" แล้ว",
         size=13, color=COLORS['muted'])

# Bad card
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(2),
                           Inches(6), Inches(4.5))
shape.fill.solid(); shape.fill.fore_color.rgb = COLORS['panel']
shape.line.color.rgb = COLORS['red']; shape.line.width = Pt(2)
tf = shape.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "❌ ใช้แค่ราคา (close)"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['red']
for line in ["", "2050.20", "2051.05", "2050.85", "2051.40", "...", "",
             "Model ต้องคิดเองว่า trend ไปไหน",
             "→ เรียนยาก ต้องใช้ data เยอะมาก"]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(13); rr.font.color.rgb = COLORS['text']
    if line.startswith("→"): rr.font.color.rgb = COLORS['red']

# Good card
shape2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.8), Inches(2),
                            Inches(6), Inches(4.5))
shape2.fill.solid(); shape2.fill.fore_color.rgb = COLORS['panel']
shape2.line.color.rgb = COLORS['p2']; shape2.line.width = Pt(2)
tf2 = shape2.text_frame; tf2.margin_left = Inches(0.3); tf2.margin_top = Inches(0.2)
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "✅ ใช้ Features (35 ตัว)"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = COLORS['p2']
for line in ["", "RSI = 32 (oversold)", "EMA_fast > EMA_slow (uptrend)",
             "ATR rising (volatility ↑)", "BB_position = 0.15 (near bottom)", "...", "",
             "Indicators แปลความหมายแล้ว",
             "→ Model จับ pattern ได้ง่าย"]:
    pp = tf2.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(13); rr.font.color.rgb = COLORS['text']
    if line.startswith("→"): rr.font.color.rgb = COLORS['p2']

add_text(s, 0.5, 6.7, 12, 0.4,
         "หลัก: \"Better features beat better algorithms\" — Pedro Domingos",
         size=13, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE 5: Phase 2 Title ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0, 1.5, 13.3, 0.5, "PHASE 2",
         size=18, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 2.2, 13.3, 1.5, "🧠 TRAIN + TEST",
         size=60, bold=True, color=COLORS['p2'], align=PP_ALIGN.CENTER)
add_text(s, 0, 4, 13.3, 1.5, "🎓", size=80, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.5, 13.3, 0.6,
         "\"เทรน RL Agent → Validate ก่อน Deploy\"",
         size=20, color=COLORS['text'], align=PP_ALIGN.CENTER)

# === SLIDE 6: Phase 2 Pipeline ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🔄 Pipeline ของ Phase 2",
         size=28, bold=True, color=COLORS['p2'])

steps = [("📥", "CSV"), ("⚖️", "Relabel"), ("🧠", "Train PPO"),
         ("📊", "Backtest"), ("🔬", "Walk-Forward")]
for i, (icon, name) in enumerate(steps):
    left = 0.5 + i * 2.5
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(2),
                                Inches(2.2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['panel']
    shape.line.color.rgb = COLORS['p2']
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    p1.add_run().text = icon
    p1.runs[0].font.size = Pt(36)

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    p2.add_run().text = name
    p2.runs[0].font.size = Pt(13)
    p2.runs[0].font.bold = True
    p2.runs[0].font.color.rgb = COLORS['p2']

add_code(s, 0.5, 4.2, 12.3, 2.5, """python relabel.py training_data_v2.csv
python rl_train.py training_data_v2_relabeled.csv --steps 200000
python rl_backtest.py rl_v3 training_data_v2_relabeled.csv
python rl_walkforward.py training_data_v2_relabeled.csv --windows 5""")

# === SLIDE: What is PPO? (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🤖 PPO คืออะไร? (อธิบายง่ายๆ)",
         size=28, bold=True, color=COLORS['p2'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "PPO = Proximal Policy Optimization — Algorithm สอน Neural Network ให้ตัดสินใจดีขึ้น",
         size=13, color=COLORS['muted'])

add_card(s, 0.5, 2, 6, 4.5,
         "🎯 ทำงานยังไง",
         "1. Agent ดู state → ตัดสินใจ\n"
         "   (Buy/Sell/Hold)\n"
         "2. ลงมือทำ → เห็น reward (+/-)\n"
         "3. ปรับ NN: action ดี → เพิ่ม prob,\n"
         "   action แย่ → ลด prob\n"
         "4. ทำซ้ำ 200,000 ครั้ง → เก่งขึ้น",
         accent_color=COLORS['p2'])

add_card(s, 6.8, 2, 6, 4.5,
         "💡 \"Proximal\" = ใกล้ๆ",
         "ห้าม policy เปลี่ยนเร็วเกิน ±20%\nต่อ update\n\n"
         "❌ ก่อน PPO: Update เร็ว → ลืมของเก่า\n"
         "✅ PPO: Clip 0.2 → ค่อยๆ ปรับ stable\n\n"
         "= เปรียบเหมือนสอนสุนัข\n"
         "ค่อยๆ ฝึก ไม่เปลี่ยนวิธีทุกวัน",
         accent_color=COLORS['p2'])

add_text(s, 0.5, 6.7, 12, 0.4,
         "PPO = Standard ของ RL Trading ใน 2024+ (default ของ stable-baselines3)",
         size=13, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE 7: Walk-Forward ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🔬 Walk-Forward Validation Results",
         size=28, bold=True, color=COLORS['p2'])

results = [
    ("Window 1", "2018-2020", "1.07", "+6.7%", "✓", COLORS['p2']),
    ("Window 2", "2020-2021", "1.28", "+47.4%", "✓", COLORS['p2']),
    ("Window 3", "2021-2022", "1.10", "+12.3%", "✓", COLORS['p2']),
    ("Window 4", "2022-2024", "0.97", "-5.4%", "✗", COLORS['red']),
    ("Window 5", "2024-2026", "0.86", "-18.6%", "✗", COLORS['red']),
]
y0 = 1.6
for i, (win, period, pf, ret, status, color) in enumerate(results):
    y = y0 + i * 0.65
    for j, (val, w) in enumerate([(win, 1.5), (period, 2.5), (pf, 1.5),
                                    (ret, 1.8), (status, 0.8)]):
        x = 1 + sum([1.5, 2.5, 1.5, 1.8][:j])
        col = color if j == 4 else COLORS['text']
        add_text(s, x, y, w, 0.5, val, size=14, color=col)

add_text(s, 0.5, 5.5, 12, 0.5,
         "Verdict: 3/5 ผ่าน → Edge เคยมี แต่กำลัง drift → ต้อง Fine-tune",
         size=16, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE 8: Phase 3 Title ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0, 1.5, 13.3, 0.5, "PHASE 3",
         size=18, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 2.2, 13.3, 1.5, "🌐 MT5 BRIDGE",
         size=60, bold=True, color=COLORS['p3'], align=PP_ALIGN.CENTER)
add_text(s, 0, 4, 13.3, 1.5, "🚀", size=80, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.5, 13.3, 0.6,
         "\"Demo → Small Live → Scale Up\"",
         size=20, color=COLORS['text'], align=PP_ALIGN.CENTER)

# === SLIDE 9: Architecture ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🏗️ Python ↔ MT5 Architecture",
         size=28, bold=True, color=COLORS['p3'])

add_card(s, 1, 2, 5, 4,
         "🐍 Python (Brain)",
         "• live_trader.py — main loop\n"
         "• features.py — calc 35 features\n"
         "• rl_v3.zip — trained model\n"
         "• live_trades.db — SQLite log",
         accent_color=COLORS['p3'])

add_card(s, 7.3, 2, 5, 4,
         "🏦 MT5 (Hands)",
         "• MT5 Terminal (broker)\n"
         "• Live tick data\n"
         "• Order execution\n"
         "• Demo → Live account",
         accent_color=COLORS['p2'])

add_text(s, 0.5, 6.5, 12, 0.5,
         "🔌 Bridge: pip install MetaTrader5 (Official API — ไม่ต้องใช้ socket)",
         size=14, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE 10: Risk Management ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🛡️ Risk Management — 3 Layers",
         size=28, bold=True, color=COLORS['p3'])

risks = [
    ("🔴 L1", "Hard Stop (Critical)",
     "Equity < 85% → ปิดทุก position + หยุดเทรด", COLORS['red']),
    ("🟡 L2", "Soft Pause (Warning)",
     "Rolling PF < 0.9 (50 trades) → หยุด 24 ชม + trigger fine-tune",
     COLORS['yellow']),
    ("🟢 L3", "Position Sizing (Always)",
     "Risk 1% per trade · Max 3 positions · Dynamic lot",
     COLORS['p2']),
]
for i, (level, name, desc, color) in enumerate(risks):
    y = 1.7 + i * 1.6
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(y),
                                Inches(12.3), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['panel']
    shape.line.color.rgb = color
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.2)
    p1 = tf.paragraphs[0]
    r = p1.add_run(); r.text = f"{level}  "
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = color
    r2 = p1.add_run(); r2.text = name
    r2.font.size = Pt(20); r2.font.bold = True; r2.font.color.rgb = COLORS['text']

    p2 = tf.add_paragraph()
    r3 = p2.add_run(); r3.text = desc
    r3.font.size = Pt(13); r3.font.color.rgb = COLORS['muted']

# === SLIDE: Concept Drift (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8,
         "💧 Concept Drift — ทำไม Model เก่าใช้ไม่ได้นาน?",
         size=26, bold=True, color=COLORS['p4'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "Model เทรนวันนี้ดี ปีหน้าอาจไม่ได้แล้ว — เพราะตลาดเปลี่ยน",
         size=13, color=COLORS['muted'])

add_card(s, 0.5, 2, 6, 4,
         "📉 ทำไมเกิด Drift?",
         "• Algorithm trader เพิ่ม → arbitrage หาย\n"
         "• นโยบายเศรษฐกิจเปลี่ยน (FED, ECB)\n"
         "• เหตุการณ์ใหม่ (COVID, สงคราม, AI boom)\n"
         "• Liquidity เปลี่ยน\n"
         "• Pattern ที่เคย work → คนใช้กันหมด → ไม่ work",
         accent_color=COLORS['red'])

add_card(s, 6.8, 2, 6, 4,
         "✅ ทางแก้",
         "• ตรวจจับ: Walk-forward validation\n"
         "• วัดผล live: rolling PF (50 trades)\n"
         "• Update: Fine-tune ทุก 3 เดือน\n"
         "• Reset: Retrain scratch ทุก 1 ปี",
         accent_color=COLORS['p2'])

add_text(s, 0.5, 6.3, 12, 0.6,
         "\"ตลาดไม่เคยหยุดเปลี่ยน → Model ก็ต้องไม่หยุดเรียน\"",
         size=18, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# === SLIDE 11: Phase 4 Title ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0, 1.5, 13.3, 0.5, "PHASE 4",
         size=18, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 2.2, 13.3, 1.5, "🔄 FINE-TUNE",
         size=60, bold=True, color=COLORS['p4'], align=PP_ALIGN.CENTER)
add_text(s, 0, 4, 13.3, 1.5, "⚙️", size=80, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.5, 13.3, 0.6,
         "\"Adapt or Die — Edge ไม่ Permanent\"",
         size=20, color=COLORS['text'], align=PP_ALIGN.CENTER)

# === SLIDE 12: 3 Levels ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🎚️ 3 Levels of Automation",
         size=28, bold=True, color=COLORS['p4'])

levels = [
    ("👋 Level 1: MANUAL",
     "เริ่มต้น 2-3 รอบแรก\n• รัน command เอง\n• ดูผลทุก step\n• เข้าใจทุกอย่าง"),
    ("⏰ Level 2: SCHEDULED",
     "หลังคุ้นเคย\n• Task Scheduler / cron\n• ทุก 90 วันอัตโนมัติ\n• ไม่ต้องจำ"),
    ("🤖 Level 3: TRIGGER-BASED",
     "Smart auto\n• Drift detection\n• Auto เมื่อ PF drop\n• Cooldown + gates"),
]
for i, (title, content) in enumerate(levels):
    add_card(s, 0.5 + i * 4.2, 1.7, 4, 4.5,
             title, content, accent_color=COLORS['p4'])

# === SLIDE 13: Yearly Cycle ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "📅 ตัวอย่าง 1 ปี Production",
         size=28, bold=True, color=COLORS['p4'])

quarters = [
    ("Q1", "Trade with rl_v3", "—", "Live trades", COLORS['text']),
    ("Q2", "Fine-tune", "50k", "rl_v4", COLORS['p2']),
    ("Q3", "Fine-tune", "50k", "rl_v5", COLORS['p2']),
    ("Q4", "Fine-tune", "50k", "rl_v6", COLORS['p2']),
    ("EOY", "RETRAIN SCRATCH", "200k", "rl_v7 (clean reset)", COLORS['yellow']),
]
y0 = 1.6
add_text(s, 0.5, y0, 1.5, 0.4, "Quarter",
         size=14, bold=True, color=COLORS['yellow'])
add_text(s, 2.5, y0, 3.5, 0.4, "Activity",
         size=14, bold=True, color=COLORS['yellow'])
add_text(s, 6.5, y0, 1.5, 0.4, "Steps",
         size=14, bold=True, color=COLORS['yellow'])
add_text(s, 8.5, y0, 4, 0.4, "Output",
         size=14, bold=True, color=COLORS['yellow'])

for i, (q, act, steps, out, color) in enumerate(quarters):
    y = y0 + 0.7 + i * 0.7
    add_text(s, 0.5, y, 1.5, 0.5, q, size=16, bold=True, color=COLORS['text'])
    add_text(s, 2.5, y, 3.5, 0.5, act, size=14, color=color)
    add_text(s, 6.5, y, 1.5, 0.5, steps, size=14, color=COLORS['text'])
    add_text(s, 8.5, y, 4, 0.5, out, size=14, color=COLORS['text'])

# === SLIDE: 5 Common Mistakes (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🚨 5 ข้อผิดพลาดที่คนทำบ่อย",
         size=28, bold=True, color=COLORS['red'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "เจอกันบ่อย — ระวังไว้จะประหยัดเวลาได้เยอะ",
         size=13, color=COLORS['muted'])

mistakes = [
    ("1️⃣ Data Leakage",
     "อาการ: WR 100% ดูเทพ\nสาเหตุ: future_return ปนใน features\n✅ Drop leaky cols ก่อนเทรน"),
    ("2️⃣ Class Imbalance",
     "อาการ: Trades = 0\nสาเหตุ: FLAT 73% → model ขี้เกียจ\n✅ Relabel quantile 33/33/33"),
    ("3️⃣ Buy-and-Hold Trap",
     "อาการ: เปิด long ครั้งเดียว ถือ 3 ปี\nสาเหตุ: MTM reward\n✅ Realized + max_hold_bars"),
    ("4️⃣ Skip Walk-Forward",
     "อาการ: Backtest ดี → live พัง\nสาเหตุ: Overfit ที่ซ่อนอยู่\n✅ WF 5 windows ทุกครั้ง"),
    ("5️⃣ Forgot Maintenance",
     "อาการ: Deploy แล้วทิ้ง → กำไรหาย\nสาเหตุ: Concept drift\n✅ Fine-tune ทุก 3 เดือน"),
    ("💡 จำไว้",
     "คนเก่ง = ตรวจเจอเร็ว + แก้ได้\nไม่ใช่ \"ไม่เคยทำผิด\"\nทุกข้อข้างบนเป็นเรื่องปกติ!"),
]
for i, (title, content) in enumerate(mistakes):
    row = i // 3
    col = i % 3
    color = COLORS['red'] if i < 5 else COLORS['p2']
    add_card(s, 0.5 + col * 4.2, 1.9 + row * 2.5, 4, 2.3,
             title, content, accent_color=color)

# === SLIDE 14: DO/DON'T ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "🎯 Mindset สำคัญ",
         size=28, bold=True, color=COLORS['yellow'])

# DO
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(1.5),
                           Inches(6), Inches(5.5))
shape.fill.solid(); shape.fill.fore_color.rgb = COLORS['panel']
shape.line.color.rgb = COLORS['p2']; shape.line.width = Pt(3)
tf = shape.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2); tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "✅ DO"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = COLORS['p2']
do_items = [
    "เริ่มเล็ก, validate ทุกขั้นตอน",
    "ใช้ demo ก่อน live เสมอ",
    "Walk-forward ก่อน deploy",
    "Backup model เก่าก่อน update",
    "Risk management = ห้ามข้าม",
    "Log ทุก trade + state",
    "ยอมรับว่า edge มีวันหมด",
]
for item in do_items:
    pp = tf.add_paragraph()
    rr = pp.add_run(); rr.text = "• " + item
    rr.font.size = Pt(14); rr.font.color.rgb = COLORS['text']

# DON'T
shape2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.8), Inches(1.5),
                            Inches(6), Inches(5.5))
shape2.fill.solid(); shape2.fill.fore_color.rgb = COLORS['panel']
shape2.line.color.rgb = COLORS['red']; shape2.line.width = Pt(3)
tf2 = shape2.text_frame
tf2.margin_left = Inches(0.3); tf2.margin_top = Inches(0.2); tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "❌ DON'T"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = COLORS['red']
dont_items = [
    "ห้าม deploy โดยไม่ผ่าน WF",
    "ห้าม auto-deploy โดยไม่มี gate",
    "ห้าม fine-tune ถี่เกิน",
    "ห้ามใช้ leverage สูงตอนเริ่ม",
    "ห้ามคาดหวังกำไร 100%",
    "ห้าม fine-tune บน new only",
    "ห้าม ignore drift signals",
]
for item in dont_items:
    pp = tf2.add_paragraph()
    rr = pp.add_run(); rr.text = "• " + item
    rr.font.size = Pt(14); rr.font.color.rgb = COLORS['text']

# === SLIDE: Glossary (NEW) ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0.5, 0.4, 12, 0.8, "📖 Glossary — คำศัพท์สำคัญ",
         size=28, bold=True, color=COLORS['yellow'])
add_text(s, 0.5, 1.2, 12, 0.4,
         "คำศัพท์ที่ใช้ตลอดคอร์ส — เก็บไว้อ้างอิง",
         size=13, color=COLORS['muted'])

glossary = [
    ("📊 Backtest", "ทดสอบ strategy บนข้อมูลในอดีต"),
    ("⚖️ Walk-Forward", "Backtest แบบเลื่อนหน้าต่าง พิสูจน์ robust"),
    ("📐 Profit Factor", "กำไร÷ขาดทุน · >1.0 = กำไร"),
    ("📉 Drawdown", "% ที่ equity ตกจากจุดสูงสุด"),
    ("🎯 Win Rate", "% trade ที่กำไร · >50% = มี edge"),
    ("📊 Sharpe Ratio", "Risk-adjusted return · >1.0 ดี"),
    ("🔄 Fine-tune", "Train ต่อจาก model เก่าด้วย data ใหม่"),
    ("💧 Concept Drift", "ตลาดเปลี่ยน → model เก่าใช้ไม่ได้"),
    ("🌍 Environment", "\"โลก\" ที่ agent อยู่ใน RL"),
    ("🧠 Policy", "\"กฎ\" ของ agent · state → action"),
    ("📊 State", "ข้อมูลที่ agent เห็น ณ เวลานั้น"),
    ("⚡ Action", "สิ่งที่ agent ทำได้ · Hold/Buy/Sell/Close"),
]
for i, (term, defn) in enumerate(glossary):
    row = i // 4
    col = i % 4
    add_card(s, 0.4 + col * 3.2, 1.9 + row * 1.7, 3.0, 1.5,
             term, defn, accent_color=COLORS['yellow'])

# === SLIDE 15: Closing ===
s = prs.slides.add_slide(blank)
set_bg(s, COLORS['bg'])
add_text(s, 0, 2, 13.3, 1.5, "🎓 ขอบคุณครับ!",
         size=60, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)
add_text(s, 0, 3.8, 13.3, 0.6,
         "\"ML Trading ไม่ใช่ magic — เป็น engineering\"",
         size=20, color=COLORS['text'], align=PP_ALIGN.CENTER)
add_text(s, 0, 4.7, 13.3, 0.5,
         "80% ของงาน = data prep + monitoring + maintenance",
         size=16, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 5.2, 13.3, 0.5,
         "20% = model training",
         size=16, color=COLORS['muted'], align=PP_ALIGN.CENTER)
add_text(s, 0, 6, 13.3, 1, "🚀 📊 🎯",
         size=60, align=PP_ALIGN.CENTER)

# Save
output = "rl_course_slides.pptx"
prs.save(output)
print(f"[OK] Saved: {output}")
print(f"     Slides: {len(prs.slides)}")
print(f"\nเปิดด้วย PowerPoint หรือ LibreOffice Impress ได้เลย")
