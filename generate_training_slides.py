"""
Generate PowerPoint slides explaining HOW RL TRAINING WORKS
(features → weights → backprop → reward signal)

Output: rl_how_training_works.pptx
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─────── colors ───────
COLORS = {
    'bg':     RGBColor(0x0a, 0x0e, 0x14),
    'panel':  RGBColor(0x18, 0x18, 0x1b),
    'panel2': RGBColor(0x27, 0x27, 0x2a),
    'text':   RGBColor(0xe4, 0xe4, 0xe7),
    'muted':  RGBColor(0x71, 0x71, 0x7a),
    'green':  RGBColor(0x34, 0xd3, 0x99),
    'red':    RGBColor(0xef, 0x44, 0x44),
    'yellow': RGBColor(0xfb, 0xbf, 0x24),
    'blue':   RGBColor(0x60, 0xa5, 0xfa),
    'purple': RGBColor(0xa7, 0x8b, 0xfa),
    'orange': RGBColor(0xfb, 0x92, 0x3c),
    'pink':   RGBColor(0xf4, 0x72, 0xb6),
    'white':  RGBColor(0xff, 0xff, 0xff),
}


# ─────── helpers ───────
def set_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or COLORS['bg']


def add_text(slide, left, top, width, height, text,
              size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
              font_name=None):
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return tx


def add_card(slide, left, top, width, height, accent_color=None,
              fill_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or COLORS['panel']
    if accent_color:
        shape.line.color.rgb = accent_color
        shape.line.width = Pt(2.5)
    else:
        shape.line.color.rgb = COLORS['muted']
        shape.line.width = Pt(0.5)
    return shape


def add_arrow(slide, left, top, width, height, color=None):
    """Right-pointing arrow"""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color or COLORS['blue']
    s.line.fill.background()
    return s


def add_down_arrow(slide, left, top, width, height, color=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color or COLORS['blue']
    s.line.fill.background()
    return s


def add_code_box(slide, left, top, width, height, lines,
                  color=COLORS['green']):
    """Black panel with monospace code text"""
    card = add_card(slide, left, top, width, height,
                     accent_color=color, fill_color=COLORS['bg'])
    tx = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(height - 0.2))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, (text, c) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.name = 'Consolas'
        run.font.color.rgb = c or COLORS['text']


# ============================================================
# Build Presentation
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 1, 1.5, 11.3, 1.5, "🧠 How Training Works",
         size=56, bold=True, color=COLORS['blue'], align=PP_ALIGN.CENTER)
add_text(s, 1, 2.9, 11.3, 0.7,
         "เกิดอะไรขึ้นตอน Train? — Features, Weights, Reward",
         size=22, color=COLORS['text'], align=PP_ALIGN.CENTER)
add_text(s, 1, 3.7, 11.3, 0.5,
         "เห็นภาพชัด ตั้งแต่ Forward Pass ถึง Backprop",
         size=16, color=COLORS['muted'], align=PP_ALIGN.CENTER)

add_text(s, 1, 4.7, 11.3, 0.8, "📊  ⚙️  📈",
         size=58, align=PP_ALIGN.CENTER)

add_text(s, 1, 6.2, 11.3, 0.4,
         "Visual Guide · 13 slides",
         size=12, color=COLORS['muted'], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: BIG PICTURE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🌐 ภาพรวม: Train คืออะไร?",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "Agent = นักเรียน · Environment = สนามฝึก · Reward = คะแนน",
         size=16, color=COLORS['muted'])

# Loop diagram
add_card(s, 0.5, 2.2, 3.8, 1.5, accent_color=COLORS['blue'])
add_text(s, 0.7, 2.4, 3.5, 0.4, "🤖  AGENT",
         size=18, bold=True, color=COLORS['blue'])
add_text(s, 0.7, 2.9, 3.5, 0.4, "Neural Network",
         size=13, color=COLORS['text'])
add_text(s, 0.7, 3.25, 3.5, 0.4, "(125,000 weights)",
         size=12, color=COLORS['muted'], font_name='Consolas')

add_arrow(s, 4.5, 2.7, 1.0, 0.5, color=COLORS['green'])
add_text(s, 4.4, 2.2, 1.2, 0.4, "action", size=11,
         color=COLORS['green'], align=PP_ALIGN.CENTER)

add_card(s, 5.7, 2.2, 3.8, 1.5, accent_color=COLORS['orange'])
add_text(s, 5.9, 2.4, 3.5, 0.4, "🌍  ENVIRONMENT",
         size=18, bold=True, color=COLORS['orange'])
add_text(s, 5.9, 2.9, 3.5, 0.4, "TradingEnv",
         size=13, color=COLORS['text'])
add_text(s, 5.9, 3.25, 3.5, 0.4, "(price + features)",
         size=12, color=COLORS['muted'], font_name='Consolas')

add_arrow(s, 9.7, 2.7, 1.0, 0.5, color=COLORS['yellow'])
add_text(s, 9.6, 2.2, 1.2, 0.4, "reward", size=11,
         color=COLORS['yellow'], align=PP_ALIGN.CENTER)

add_card(s, 10.9, 2.2, 1.9, 1.5, accent_color=COLORS['purple'])
add_text(s, 11.0, 2.6, 1.7, 0.4, "📈",
         size=24, align=PP_ALIGN.CENTER)
add_text(s, 11.0, 3.1, 1.7, 0.4, "Update",
         size=14, bold=True, color=COLORS['purple'],
         align=PP_ALIGN.CENTER)

# Bottom: feedback loop
add_text(s, 0.5, 4.2, 12.3, 0.4,
         "↻  ทำซ้ำ ~200,000 ครั้ง — Agent เก่งขึ้นทีละนิด",
         size=18, bold=True, color=COLORS['green'], align=PP_ALIGN.CENTER)

# Key takeaway
add_card(s, 0.5, 5.0, 12.3, 1.8, accent_color=COLORS['blue'])
add_text(s, 0.8, 5.2, 11.7, 0.5, "💡 หัวใจของ Train:",
         size=16, bold=True, color=COLORS['blue'])
add_text(s, 0.8, 5.7, 11.7, 1.0,
         "1. Agent ดู state (features ของตลาด) → ตัดสินใจ action\n"
         "2. Environment ตอบกลับด้วย reward (กำไร/ขาดทุน)\n"
         "3. ขยับ weights ของ Neural Network ให้ \"เก่งขึ้น\" ในรอบถัดไป",
         size=13, color=COLORS['text'])


# ============================================================
# SLIDE 3: STATE = FEATURES VECTOR
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📊 1. Features → State Vector",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ทุก bar ของกราฟ = ภาพถ่าย 1 รูป → แปลงเป็นตัวเลข",
         size=16, color=COLORS['muted'])

# Left: a candle bar
add_card(s, 0.5, 2.2, 3.0, 4.5, accent_color=COLORS['orange'])
add_text(s, 0.7, 2.35, 2.8, 0.4, "📈 1 Bar (H1)",
         size=14, bold=True, color=COLORS['orange'])
add_text(s, 0.7, 2.85, 2.8, 0.4, "EURUSD",
         size=11, color=COLORS['muted'], font_name='Consolas')

# Simple "candle" shape
candle = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(1.7), Inches(3.5), Inches(0.5), Inches(1.5))
candle.fill.solid()
candle.fill.fore_color.rgb = COLORS['green']
candle.line.fill.background()
# wick
wick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(1.92), Inches(3.2), Inches(0.06), Inches(2.0))
wick.fill.solid()
wick.fill.fore_color.rgb = COLORS['green']
wick.line.fill.background()

add_text(s, 0.7, 5.4, 2.8, 0.4, "OHLC + Volume",
         size=12, color=COLORS['text'], align=PP_ALIGN.CENTER)
add_text(s, 0.7, 5.85, 2.8, 0.4, "+ 30 indicators",
         size=12, color=COLORS['blue'], align=PP_ALIGN.CENTER)
add_text(s, 0.7, 6.25, 2.8, 0.4, "+ position state",
         size=12, color=COLORS['purple'], align=PP_ALIGN.CENTER)

# Arrow
add_arrow(s, 3.7, 4.2, 0.7, 0.6, color=COLORS['blue'])

# Right: vector representation
add_code_box(s, 4.6, 2.2, 8.2, 4.5, [
    ("# State vector (window=10, features=33)", COLORS['muted']),
    ("[", COLORS['text']),
    ("  rsi_14       =  0.42      <- มาตรวัดแรง", COLORS['green']),
    ("  ema_diff     = -0.18      <- ทิศทาง", COLORS['green']),
    ("  atr_norm     =  1.05      <- ความผันผวน", COLORS['green']),
    ("  hour         =  0.83      <- เวลา", COLORS['green']),
    ("  ...                       (อีก 30+ ตัว x 10 bars)", COLORS['muted']),
    ("", COLORS['text']),
    ("  position       =  1       <- ถือ long", COLORS['purple']),
    ("  unrealized_pnl =  0.005   <- กำไรลอย 0.5%", COLORS['purple']),
    ("  bars_held      =  0.07    <- ถือมา 7 bars", COLORS['purple']),
    ("]", COLORS['text']),
    ("", COLORS['text']),
    ("  shape = (333,)   <- ป้อนเข้า Neural Network", COLORS['yellow']),
])


# ============================================================
# SLIDE 4: NEURAL NETWORK STRUCTURE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "⚙️ 2. Neural Network = Layers + Weights",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "Network = หลาย layer ต่อกัน — แต่ละ layer มี weights นับหมื่น",
         size=16, color=COLORS['muted'])

# Layers with arrows
layers = [
    ("Input",    "333",  COLORS['blue'],   "state vector"),
    ("Layer 1",  "256",  COLORS['green'],  "≈ 85K weights"),
    ("Layer 2",  "128",  COLORS['green'],  "≈ 33K weights"),
    ("Layer 3",  "64",   COLORS['green'],  "≈ 8K weights"),
    ("Output",   "4",    COLORS['orange'], "Hold/Buy/Sell/Close"),
]

x_pos = 0.5
for i, (name, neurons, col, sub) in enumerate(layers):
    add_card(s, x_pos, 2.4, 2.0, 2.5, accent_color=col)
    add_text(s, x_pos + 0.1, 2.55, 1.8, 0.4, name,
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x_pos + 0.1, 3.1, 1.8, 0.6, neurons,
             size=32, bold=True, color=COLORS['text'],
             align=PP_ALIGN.CENTER, font_name='Consolas')
    add_text(s, x_pos + 0.1, 3.95, 1.8, 0.4, "neurons",
             size=11, color=COLORS['muted'], align=PP_ALIGN.CENTER)
    add_text(s, x_pos + 0.1, 4.4, 1.8, 0.4, sub,
             size=10, color=COLORS['yellow'], align=PP_ALIGN.CENTER,
             font_name='Consolas')

    if i < len(layers) - 1:
        add_arrow(s, x_pos + 2.05, 3.4, 0.4, 0.4, color=COLORS['muted'])

    x_pos += 2.45

# Total weights box
add_card(s, 0.5, 5.4, 12.3, 1.5, accent_color=COLORS['purple'])
add_text(s, 0.8, 5.55, 11.7, 0.5, "🧮 Total Weights",
         size=16, bold=True, color=COLORS['purple'])
add_text(s, 0.8, 6.05, 11.7, 0.4,
         "≈ 125,000 ตัวที่ต้อง \"จูน\" — เริ่มจาก random",
         size=14, color=COLORS['text'])
add_text(s, 0.8, 6.45, 11.7, 0.4,
         "เป้าหมาย Train: หาค่า weights ทั้ง 125K ที่ทำให้ตัดสินใจได้ดีที่สุด",
         size=12, color=COLORS['muted'])


# ============================================================
# SLIDE 5: WHERE DOES 125,000 COME FROM?
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🧮 ทำไมถึงเป็น 125,000 weights?",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "weights ระหว่าง 2 layer = (input neurons) × (output neurons) + (bias)",
         size=15, color=COLORS['muted'])

# Left: input vector calc
add_card(s, 0.5, 2.1, 5.8, 1.5, accent_color=COLORS['blue'])
add_text(s, 0.7, 2.25, 5.4, 0.4, "📥 Input Vector",
         size=14, bold=True, color=COLORS['blue'])
add_text(s, 0.7, 2.7, 5.4, 0.4,
         "state = window × features + 3",
         size=12, color=COLORS['text'], font_name='Consolas')
add_text(s, 0.7, 3.1, 5.4, 0.4,
         "      = 10 × 33 + 3  =  333 ตัว",
         size=13, bold=True, color=COLORS['green'], font_name='Consolas')

# Right: formula card
add_card(s, 6.4, 2.1, 6.4, 1.5, accent_color=COLORS['purple'])
add_text(s, 6.6, 2.25, 6.0, 0.4, "🔢 Bias คืออะไร?",
         size=14, bold=True, color=COLORS['purple'])
add_text(s, 6.6, 2.7, 6.0, 0.4,
         "ค่าคงที่บวกต่อ neuron 1 ตัว",
         size=12, color=COLORS['text'])
add_text(s, 6.6, 3.1, 6.0, 0.4,
         "เหมือน 'b' ในสมการ y = mx + b",
         size=11, color=COLORS['muted'], font_name='Consolas')

# Math table — header
add_card(s, 0.5, 3.8, 12.3, 2.7, accent_color=COLORS['orange'])
add_text(s, 0.7, 3.95, 11.9, 0.4, "📐 คำนวณทีละ layer:",
         size=14, bold=True, color=COLORS['orange'])

# Code box with calculation
add_code_box(s, 0.7, 4.4, 11.9, 2.0, [
    ("Layer 1   Input(333) -> Hidden(256):  333 × 256 + 256 =  85,504", COLORS['text']),
    ("Layer 2   Hidden(256) -> Hidden(128): 256 × 128 + 128 =  32,896", COLORS['text']),
    ("Layer 3   Hidden(128) -> Hidden(64):  128 ×  64 +  64 =   8,256", COLORS['text']),
    ("Actor     Hidden(64) -> Output(4):     64 ×   4 +   4 =     260   <- Hold/Buy/Sell/Close", COLORS['green']),
    ("Critic    Hidden(64) -> Value(1):      64 ×   1 +   1 =      65   <- ทาย state value", COLORS['pink']),
    ("                                       ─────────────────────────", COLORS['muted']),
    ("                                       TOTAL          ≈ 126,981   ≈ 125K ✅", COLORS['yellow']),
])

# Bottom: scale table
add_card(s, 0.5, 6.6, 12.3, 0.8, accent_color=COLORS['blue'],
          fill_color=COLORS['panel2'])
add_text(s, 0.7, 6.7, 11.9, 0.4,
         "💡 Network ขึ้นกับ window:  window<20 → 50K · window 20–49 → 125K · window≥50 → 350K",
         size=12, color=COLORS['text'], font_name='Consolas')
add_text(s, 0.7, 7.05, 11.9, 0.4,
         "(auto-scale ใน rl_train.py — ใหญ่ไป overfit, เล็กไป underfit)",
         size=11, color=COLORS['muted'])


# ============================================================
# SLIDE 6: FORWARD PASS
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "➡️ 3. Forward Pass (state → action)",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "เอา features × weights ทีละ layer → ออกมาเป็น probability ของแต่ละ action",
         size=15, color=COLORS['muted'])

# Math flow
add_code_box(s, 0.5, 2.2, 12.3, 3.6, [
    ("state (333,)", COLORS['blue']),
    ("    ↓", COLORS['muted']),
    ("    × W1 + b1   →  ReLU  →  activations1 (256,)   <- 'รู้สึก' กับ pattern", COLORS['text']),
    ("    ↓", COLORS['muted']),
    ("    × W2 + b2   →  ReLU  →  activations2 (128,)   <- กรองข้อมูลขั้นที่ 2", COLORS['text']),
    ("    ↓", COLORS['muted']),
    ("    × W3 + b3   →  ReLU  →  activations3 (64,)    <- รวบ pattern เข้มข้น", COLORS['text']),
    ("    ↓", COLORS['muted']),
    ("    × W_out + b_out →  Softmax  →  probability vector", COLORS['text']),
    ("", COLORS['text']),
    ("OUTPUT:  [0.10,  0.60,  0.20,  0.10]", COLORS['orange']),
    ("          Hold   Buy   Sell  Close", COLORS['muted']),
    ("                  ↑", COLORS['green']),
    ("              Agent เลือก: BUY  (probability สูงสุด)", COLORS['green']),
])

# Note
add_card(s, 0.5, 6.0, 12.3, 1.0, accent_color=COLORS['yellow'],
          fill_color=COLORS['panel2'])
add_text(s, 0.8, 6.15, 11.7, 0.4, "💡 ตอนเริ่ม weights = random",
         size=14, bold=True, color=COLORS['yellow'])
add_text(s, 0.8, 6.55, 11.7, 0.4,
         "→ output ก็มั่ว ก็เหมือนคนเริ่มเล่นโป๊กเกอร์ — ต้องเล่นเยอะๆ ถึงจะเก่ง",
         size=12, color=COLORS['text'])


# ============================================================
# SLIDE 7: REWARD SIGNAL
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🎯 4. Reward = สัญญาณ \"เลือกถูกหรือผิด?\"",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "Environment คำนวณ reward ทุก step — อันนี้คือคะแนนที่ guide การเรียน",
         size=15, color=COLORS['muted'])

# Trade timeline
add_code_box(s, 0.5, 2.2, 12.3, 4.4, [
    ("Trade timeline (Buy → Hold × 4 → Close):", COLORS['blue']),
    ("", COLORS['text']),
    ("t=0  Agent: BUY at 1.0850", COLORS['text']),
    ("     Env:   เปิด long, cost = -0.0003", COLORS['red']),
    ("", COLORS['text']),
    ("t=1  Agent: HOLD", COLORS['text']),
    ("     Env:   price 1.0855  (+0.046%)", COLORS['muted']),
    ("            reward += 0.00046  <- unrealized delta (เราเพิ่มใหม่!)", COLORS['green']),
    ("", COLORS['text']),
    ("t=2  Agent: HOLD", COLORS['text']),
    ("     Env:   price 1.0860  (+0.092%)", COLORS['muted']),
    ("            reward += 0.00046", COLORS['green']),
    ("", COLORS['text']),
    ("t=5  Agent: CLOSE at 1.0870", COLORS['text']),
    ("     Env:   pnl = +0.184%", COLORS['muted']),
    ("            reward = +1.84    <- Big chunk × 10!", COLORS['green']),
    ("                   + 0.005    <- profit bonus", COLORS['green']),
    ("                   - 0.0001   <- commission", COLORS['red']),
    ("            -------- Total: ~+1.85 ----------", COLORS['yellow']),
])


# ============================================================
# SLIDE 8: TRAJECTORY & ADVANTAGE
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🗂️ 5. Rollout Buffer + Advantage",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "PPO เก็บ trajectories 2,048 transitions → คำนวณ \"advantage\" ของแต่ละ action",
         size=15, color=COLORS['muted'])

# Buffer table
add_code_box(s, 0.5, 2.2, 7.0, 4.6, [
    ("Rollout Buffer (n_steps = 2048):", COLORS['blue']),
    ("", COLORS['text']),
    ("┌────┬───────┬────────┬─────────┐", COLORS['muted']),
    ("│ t  │ state │ action │ reward  │", COLORS['muted']),
    ("├────┼───────┼────────┼─────────┤", COLORS['muted']),
    ("│ 0  │ s_0   │ Buy    │ -0.0003 │", COLORS['text']),
    ("│ 1  │ s_1   │ Hold   │ +0.0005 │", COLORS['text']),
    ("│ 2  │ s_2   │ Hold   │ +0.0005 │", COLORS['text']),
    ("│ 3  │ s_3   │ Hold   │ +0.0005 │", COLORS['text']),
    ("│ 4  │ s_4   │ Hold   │ +0.0005 │", COLORS['text']),
    ("│ 5  │ s_5   │ Close  │ +1.85   │ <-- big!", COLORS['green']),
    ("│ 6  │ s_6   │ Hold   │ -0.0001 │", COLORS['text']),
    ("│... │ ...   │ ...    │ ...     │", COLORS['muted']),
    ("└────┴───────┴────────┴─────────┘", COLORS['muted']),
    ("    × 2048 rows                  ", COLORS['muted']),
])

# Right: Advantage explanation
add_card(s, 7.7, 2.2, 5.1, 4.6, accent_color=COLORS['purple'])
add_text(s, 7.9, 2.35, 4.7, 0.5, "📐 Advantage",
         size=18, bold=True, color=COLORS['purple'])
add_text(s, 7.9, 2.85, 4.7, 0.4,
         "= action นี้ดีกว่าค่าเฉลี่ยแค่ไหน?",
         size=12, color=COLORS['muted'])

add_text(s, 7.9, 3.4, 4.7, 0.4,
         "ตัวอย่าง:",
         size=13, bold=True, color=COLORS['text'])

add_text(s, 7.9, 3.85, 4.7, 0.4,
         "• Buy ที่ t=0 → ได้ +1.85 ตอน t=5",
         size=11, color=COLORS['text'])
add_text(s, 7.9, 4.25, 4.7, 0.4,
         "  Adv = +1.5  ✅ ดีมาก!",
         size=11, color=COLORS['green'], font_name='Consolas')

add_text(s, 7.9, 4.85, 4.7, 0.4,
         "• Hold ที่ t=2 → trade ติดอีก 3 bars",
         size=11, color=COLORS['text'])
add_text(s, 7.9, 5.25, 4.7, 0.4,
         "  Adv = -0.05  ❌ ไม่ดี",
         size=11, color=COLORS['red'], font_name='Consolas')

add_text(s, 7.9, 5.85, 4.7, 0.7,
         "→ Adv > 0 = push weights ให้\n"
         "   เลือก action นี้บ่อยขึ้น",
         size=12, color=COLORS['yellow'])


# ============================================================
# SLIDE 9: GRADIENT DESCENT
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "⬇️ 6. Gradient Descent — \"ก้าวลงเขา\"",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ขยับ weights ทีละนิด ในทิศที่ทำให้ \"loss\" ลดลง",
         size=15, color=COLORS['muted'])

# Formula
add_code_box(s, 0.5, 2.2, 12.3, 1.6, [
    ("Update rule:", COLORS['blue']),
    ("", COLORS['text']),
    ("    W_new  =  W_old  -  learning_rate  ×  gradient", COLORS['green']),
    ("                              ↑                  ↑", COLORS['muted']),
    ("                            3e-4         ทิศที่ทำให้ loss ลดลง", COLORS['muted']),
])

# Loss landscape ASCII
add_code_box(s, 0.5, 4.0, 7.5, 3.0, [
    ("Loss landscape:", COLORS['blue']),
    ("", COLORS['text']),
    ("  loss", COLORS['muted']),
    ("    ▲", COLORS['muted']),
    ("    │ ╲                ╱", COLORS['red']),
    ("    │  ╲              ╱   <- weights แย่", COLORS['text']),
    ("    │   ╲    ⬇        ╱", COLORS['yellow']),
    ("    │    ╲  ก้าว     ╱     <- ก้าวลงเขา", COLORS['text']),
    ("    │     ╲ลงเขา    ╱", COLORS['yellow']),
    ("    │      ╲────────╱", COLORS['green']),
    ("    │       ⭐ minimum    <- weights ดีที่สุด", COLORS['green']),
    ("    └──────────────────► weight value", COLORS['muted']),
])

# Right: lr too big/small
add_card(s, 8.3, 4.0, 4.5, 3.0, accent_color=COLORS['purple'])
add_text(s, 8.5, 4.15, 4.2, 0.4, "⚠️ Learning Rate",
         size=15, bold=True, color=COLORS['purple'])

add_text(s, 8.5, 4.65, 4.2, 0.4, "🚀 lr สูงไป (5e-3):",
         size=12, bold=True, color=COLORS['red'])
add_text(s, 8.5, 5.05, 4.2, 0.4,
         "ก้าวใหญ่ → กระโดดข้ามเป้า",
         size=11, color=COLORS['text'])
add_text(s, 8.5, 5.4, 4.2, 0.4,
         "→ KL spike, training ไม่นิ่ง",
         size=10, color=COLORS['muted'])

add_text(s, 8.5, 5.95, 4.2, 0.4, "🐌 lr ต่ำไป (1e-6):",
         size=12, bold=True, color=COLORS['yellow'])
add_text(s, 8.5, 6.35, 4.2, 0.4,
         "ก้าวเล็ก → ช้ามาก",
         size=11, color=COLORS['text'])
add_text(s, 8.5, 6.7, 4.2, 0.4,
         "→ train ไม่ทันจบ",
         size=10, color=COLORS['muted'])


# ============================================================
# SLIDE 10: WEIGHTS ขยับ ภาพย่อย
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🔧 7. Weights ขยับยังไง?",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ทุก weight ใน 125,000 ตัว ขยับนิดเดียว → รวมกันเปลี่ยน behavior",
         size=15, color=COLORS['muted'])

# Before/After example
add_card(s, 0.5, 2.2, 6.0, 4.4, accent_color=COLORS['red'])
add_text(s, 0.7, 2.35, 5.6, 0.4, "❌ ก่อน Train (random)",
         size=15, bold=True, color=COLORS['red'])
add_code_box(s, 0.7, 2.85, 5.6, 3.6, [
    ("# State: ราคาเริ่มขึ้น + RSI ต่ำ", COLORS['muted']),
    ("# (sweet spot ของการ Buy)", COLORS['muted']),
    ("", COLORS['text']),
    ("Network output:", COLORS['text']),
    ("[0.25, 0.25, 0.25, 0.25]", COLORS['muted']),
    (" Hold   Buy  Sell Close", COLORS['muted']),
    ("", COLORS['text']),
    ("→ Agent เลือกมั่ว", COLORS['red']),
    ("  มี chance Buy ถูกแค่ 25%", COLORS['red']),
])

add_card(s, 6.8, 2.2, 6.0, 4.4, accent_color=COLORS['green'])
add_text(s, 7.0, 2.35, 5.6, 0.4, "✅ หลัง Train 200K steps",
         size=15, bold=True, color=COLORS['green'])
add_code_box(s, 7.0, 2.85, 5.6, 3.6, [
    ("# State: ราคาเริ่มขึ้น + RSI ต่ำ", COLORS['muted']),
    ("# (sweet spot ของการ Buy)", COLORS['muted']),
    ("", COLORS['text']),
    ("Network output:", COLORS['text']),
    ("[0.10, 0.78, 0.08, 0.04]", COLORS['green']),
    (" Hold   Buy  Sell Close", COLORS['muted']),
    ("", COLORS['text']),
    ("→ Agent มั่นใจเลือก Buy", COLORS['green']),
    ("  78% probability — เก่งแล้ว!", COLORS['green']),
])

# Bottom: how
add_card(s, 0.5, 6.8, 12.3, 0.6, accent_color=COLORS['blue'])
add_text(s, 0.7, 6.9, 11.9, 0.4,
         "↑ เกิดจาก: weight แต่ละตัวที่เกี่ยวกับ \"RSI ต่ำ + ราคาขึ้น\" → \"Buy\" ขยับเพิ่มขึ้นเรื่อยๆทุก epoch",
         size=12, color=COLORS['text'])


# ============================================================
# SLIDE 11: REWARD SHAPING (NEW!)
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "⭐ 8. ทำไมต้องมี Reward ระหว่าง Hold?",
         size=32, bold=True, color=COLORS['orange'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ปัญหา reward เก่า + วิธีแก้ที่เพิ่งใส่เข้าไปใน trading_env.py",
         size=15, color=COLORS['muted'])

# Before
add_card(s, 0.5, 2.2, 6.0, 4.5, accent_color=COLORS['red'])
add_text(s, 0.7, 2.35, 5.6, 0.4, "❌ ของเดิม: reward = 0 ระหว่าง Hold",
         size=14, bold=True, color=COLORS['red'])
add_code_box(s, 0.7, 2.85, 5.6, 3.7, [
    ("Buy → Hold → Hold → Hold → Hold → Close (+1.84)", COLORS['text']),
    (" │     0      0      0      0      ↑", COLORS['muted']),
    (" │                                  │", COLORS['muted']),
    (" └── 30 bars นี้ NO signal! ────────┘", COLORS['red']),
    ("", COLORS['text']),
    ("ผลลัพธ์:", COLORS['yellow']),
    ("• Agent ไม่รู้ว่า Hold ตอนไหนดี", COLORS['text']),
    ("• Close timing พัง", COLORS['text']),
    ("• ต้องพึ่ง SL/TP ช่วยปิด", COLORS['text']),
    ("", COLORS['text']),
    ("→ Pure Agent backtest:", COLORS['muted']),
    ("   PF = 0.58  (ขาดทุน)", COLORS['red']),
])

# After
add_card(s, 6.8, 2.2, 6.0, 4.5, accent_color=COLORS['green'])
add_text(s, 7.0, 2.35, 5.6, 0.4, "✅ ใหม่: incremental unrealized",
         size=14, bold=True, color=COLORS['green'])
add_code_box(s, 7.0, 2.85, 5.6, 3.7, [
    ("Buy → Hold(+0.5) → Hold(+0.5) → Hold(-0.5)", COLORS['text']),
    ("       ↑              ↑              ↑", COLORS['muted']),
    ("   ราคาขึ้น       ราคาขึ้น     ราคาลง", COLORS['muted']),
    ("   = signal+      signal+      signal-", COLORS['green']),
    ("                                ↓", COLORS['text']),
    ("                          = บอกให้ Close!", COLORS['yellow']),
    ("", COLORS['text']),
    ("ผลที่คาดว่าจะได้:", COLORS['yellow']),
    ("• เรียน close timing ได้ดีขึ้น", COLORS['text']),
    ("• ไม่ต้องพึ่ง SL/TP มาก", COLORS['text']),
    ("• Pure Agent PF น่าจะ > 1.0", COLORS['green']),
])

# Bottom callout
add_text(s, 0.5, 6.95, 12.3, 0.4,
         "💡 Gradient signal ทุก bar = weights ที่เกี่ยวกับ \"close timing\" จะขยับได้",
         size=13, bold=True, color=COLORS['blue'], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: WHOLE TRAINING JOURNEY
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "📈 9. ภาพรวมการเรียนรู้ตลอด 200K steps",
         size=32, bold=True, color=COLORS['yellow'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "ดู ep_rew_mean ตลอด training — แต่ละ phase มีอาการต่างกัน",
         size=15, color=COLORS['muted'])

phases = [
    ("0 - 50K",     "🐣", "Random Phase",
     "weights มั่ว · reward ใกล้ 0 · มี exploration เยอะ", COLORS['muted']),
    ("50K - 100K",  "📚", "Learning Phase",
     "เริ่มเห็น pattern · reward ค่อยๆขึ้น · KL/Clip ปกติ", COLORS['blue']),
    ("100K - 150K", "🚀", "Convergence Phase",
     "เก่งขึ้นชัด · reward stable · entropy ลดลง", COLORS['green']),
    ("150K - 200K", "⚠️", "Watch Out Phase",
     "อาจ overfit · ถ้า reward ตก → popup เด้ง", COLORS['orange']),
]

for i, (steps, icon, title, desc, color) in enumerate(phases):
    y = 2.2 + i * 1.2
    add_card(s, 0.5, y, 12.3, 1.0, accent_color=color)
    add_text(s, 0.7, y + 0.2, 1.0, 0.6, icon,
             size=28, align=PP_ALIGN.CENTER)
    add_text(s, 1.8, y + 0.1, 2.3, 0.4, steps,
             size=14, bold=True, color=color, font_name='Consolas')
    add_text(s, 1.8, y + 0.5, 2.3, 0.4, title,
             size=13, color=COLORS['text'])
    add_text(s, 4.2, y + 0.3, 8.4, 0.4, desc,
             size=12, color=COLORS['muted'])


# ============================================================
# SLIDE 13: SUMMARY
# ============================================================
s = prs.slides.add_slide(blank)
set_bg(s)

add_text(s, 0.5, 0.4, 12, 0.8, "🎯 สรุป — Train คืออะไร?",
         size=32, bold=True, color=COLORS['green'])

add_text(s, 0.5, 1.3, 12, 0.5,
         "1 ประโยค + 5 ขั้นตอน",
         size=15, color=COLORS['muted'])

# Big quote
add_card(s, 0.5, 2.0, 12.3, 1.4, accent_color=COLORS['yellow'],
          fill_color=COLORS['panel2'])
add_text(s, 0.8, 2.2, 11.7, 1.0,
         "\"Train = ป้อน features เข้า → ดูว่า action ที่เลือก\n"
         " ได้ reward ดีไหม → ขยับ weights ทุกตัวเล็กน้อย\n"
         " ให้ \"เลือกแบบเดิม\" บ่อยขึ้นเมื่อเจอ state คล้ายๆเดิม\"",
         size=15, bold=True, color=COLORS['yellow'], align=PP_ALIGN.CENTER)

# 5 steps
steps = [
    ("1", "📊", "Features", "ตลาด → ตัวเลข vector", COLORS['blue']),
    ("2", "⚙️", "Forward", "× weights → action prob", COLORS['purple']),
    ("3", "🎯", "Reward", "env บอกว่าดี/แย่", COLORS['orange']),
    ("4", "📐", "Advantage", "ดีกว่าค่าเฉลี่ยไหม?", COLORS['pink']),
    ("5", "⬇️", "Backprop", "ขยับ weights ลงเขา", COLORS['green']),
]

for i, (num, icon, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_card(s, x, 3.7, 2.3, 2.5, accent_color=color)
    add_text(s, x + 0.1, 3.85, 2.1, 0.4, "Step " + num,
             size=11, color=COLORS['muted'], align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 4.25, 2.1, 0.5, icon,
             size=28, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 4.95, 2.1, 0.4, title,
             size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 5.45, 2.1, 0.7, desc,
             size=11, color=COLORS['text'], align=PP_ALIGN.CENTER)

# Footer
add_text(s, 0.5, 6.6, 12.3, 0.4,
         "ทำซ้ำ ~200K ครั้ง = Network จำ pattern ตลาดได้",
         size=14, bold=True, color=COLORS['blue'], align=PP_ALIGN.CENTER)
add_text(s, 0.5, 7.0, 12.3, 0.3,
         "🧠 RL Trading Pipeline · How Training Works",
         size=10, color=COLORS['muted'], align=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output = "rl_how_training_works.pptx"
prs.save(output)
print("[ok] saved -> " + output)
print("     slides: " + str(len(prs.slides)))
